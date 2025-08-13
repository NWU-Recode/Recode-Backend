from pathlib import Path
import sys

import pytest
from fastapi import FastAPI
from fastapi.testclient import TestClient
from pptx import Presentation
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker
from sqlalchemy.pool import StaticPool

# Allow "app" imports when running tests directly
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from app.db.base import Base
from app.db.client import get_db           # <- updated import (was app.db.session)
from app.features.slide_extraction.endpoints import router as slide_router
from app.features.slide_extraction.pptx_extraction import extract_pptx_text
from app.features.slide_extraction import repository, schemas


application = FastAPI()
application.include_router(slide_router)

# In-memory SQLite for isolated tests
engine = create_engine(
    "sqlite:///:memory:",
    connect_args={"check_same_thread": False},
    poolclass=StaticPool,
)
TestingSessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base.metadata.create_all(bind=engine)


def override_get_db():
    db = TestingSessionLocal()
    try:
        yield db
        db.commit()  # commit after request handlers using the session
    finally:
        db.close()


application.dependency_overrides[get_db] = override_get_db


def _create_sample_pptx(path: Path) -> None:
    prs = Presentation()
    layout = prs.slide_layouts[0]

    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Title 1"
    slide1.placeholders[1].text = "Subtitle 1"

    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "Title 2"
    slide2.placeholders[1].text = "Subtitle 2"

    prs.save(path)


def test_extract_pptx_text(tmp_path: Path) -> None:
    pptx_file = tmp_path / "sample.pptx"
    _create_sample_pptx(pptx_file)

    result = extract_pptx_text(str(pptx_file))

    assert result == {
        1: ["Title 1", "Subtitle 1"],
        2: ["Title 2", "Subtitle 2"],
    }


def test_extract_pptx_text_invalid_file(tmp_path: Path) -> None:
    txt_file = tmp_path / "sample.txt"
    txt_file.write_text("not a pptx")

    with pytest.raises(ValueError):
        extract_pptx_text(str(txt_file))


def test_extract_endpoint(tmp_path: Path) -> None:
    pptx_file = tmp_path / "sample.pptx"
    _create_sample_pptx(pptx_file)

    client = TestClient(application)
    with pptx_file.open("rb") as f:
        response = client.post(
            "/slides/extract",
            files={
                "file": (
                    "sample.pptx",
                    f.read(),
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            },
        )

    assert response.status_code == 200
    data = response.json()
    assert data["filename"] == "sample.pptx"
    assert data["slides"] == {
        "1": ["Title 1", "Subtitle 1"],
        "2": ["Title 2", "Subtitle 2"],
    }
    assert data["id"] == 1
    assert "created_at" in data

    bad_response = client.post(
        "/slides/extract",
        files={"file": ("sample.txt", b"notpptx", "text/plain")},
    )
    assert bad_response.status_code == 400


def test_repository_multiple_inserts(tmp_path: Path) -> None:
    db = TestingSessionLocal()
    try:
        initial = len(repository.list_extractions(db))
        data1 = schemas.SlideExtractionCreate(filename="a.pptx
