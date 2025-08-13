from __future__ import annotations

from pathlib import Path
from typing import BinaryIO, Dict, List, Union

from pptx import Presentation
from pptx.exc import PackageNotFoundError


def _load_presentation(source: Union[str, Path, BinaryIO]) -> Presentation:
    """Load a presentation from a file or binary data."""
    try:
        if isinstance(source, (str, Path)):
            path = Path(source)
            if path.suffix.lower() != ".pptx":
                raise ValueError("Only .pptx files are supported")
            if not path.exists():
                raise ValueError(f"PPTX file not found: {path}")
            return Presentation(str(path))
        return Presentation(source)
    except (PackageNotFoundError, FileNotFoundError) as exc:
        raise ValueError("Invalid PPTX file") from exc


def extract_pptx_text(source: Union[str, Path, BinaryIO]) -> Dict[int, List[str]]:

    presentation = _load_presentation(source)
    slides: Dict[int, List[str]] = {}

    for index, slide in enumerate(presentation.slides, start=1):
        texts: List[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    texts.append(text)
        slides[index] = texts

    return slides


if __name__ == "__main__":
    import sys
    if len(sys.argv) != 2:
        print("Usage: python pptx_extraction.py <path-to-pptx>")
        raise SystemExit(1)
    result = extract_pptx_text(sys.argv[1])
    print(result)
