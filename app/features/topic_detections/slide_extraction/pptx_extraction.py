from __future__ import annotations

from pathlib import Path
from typing import BinaryIO, Dict, List, Union

def _import_pptx():
    try:
        from pptx import Presentation
        from pptx.exc import PackageNotFoundError
        return Presentation, PackageNotFoundError
    except Exception as e:
        raise ImportError("python-pptx is required to extract PPTX files. Install with: pip install python-pptx") from e


# Load a presentation from a file or binary data.
def _load_presentation(source: Union[str, Path, BinaryIO]):
    Presentation, PackageNotFoundError = _import_pptx()
    try:
        if isinstance(source, (str, Path)):
            path = Path(source)
            if path.suffix.lower() != ".pptx":
                raise ValueError("Only .pptx files are supported")
            if not path.exists():
                raise ValueError(f"PPTX file not found: {path}")
            return Presentation(str(path))
        return Presentation(source)
    except (PackageNotFoundError, FileNotFoundError) as exc:
        raise ValueError("Invalid PPTX file") from exc


def extract_pptx_text(source: Union[str, Path, BinaryIO]) -> Dict[int, List[str]]:
    presentation = _load_presentation(source)
    slides: Dict[int, List[str]] = {}

    for index, slide in enumerate(presentation.slides, start=1):
        texts: List[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    texts.append(text)
        slides[index] = texts

    return slides

