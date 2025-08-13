"""Utility for extracting text from PPTX files.

This script provides a simple function to read a PowerPoint presentation and
return a dictionary mapping slide numbers to lists of text found on each slide.
"""
from __future__ import annotations

from typing import Dict, List

from pptx import Presentation


def extract_pptx_text(path: str) -> Dict[int, List[str]]:
    """Extract all text from the presentation located at ``path``.

    Parameters
    ----------
    path: str
        Path to the ``.pptx`` file.

    Returns
    -------
    Dict[int, List[str]]
        Dictionary where keys are 1-based slide numbers and values are lists of
        text strings found on the corresponding slide.
    """
    presentation = Presentation(path)
    slides: Dict[int, List[str]] = {}

    for index, slide in enumerate(presentation.slides, start=1):
        texts: List[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    texts.append(text)
        slides[index] = texts

    return slides


if __name__ == "__main__":
    import sys
    if len(sys.argv) != 2:
        print("Usage: python pptx_extraction.py <path-to-pptx>")
        raise SystemExit(1)
    result = extract_pptx_text(sys.argv[1])
    print(result)
